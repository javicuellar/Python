import os
from PyPDF2 import PdfReader
from pptx import Presentation



def pdf_to_text():
    pdf_file = "05_Convertir_pdf_ppt.pdf"
    text = ""
    with open(pdf_file, 'rb') as f:
        reader = PdfReader(f)
        for page_num in range(len(reader.pages)):
            page_text = reader.pages[page_num].extract_text()
            text += page_text
    return text



def pdf_to_ppt(output_file):
    text = pdf_to_text()
    prs = Presentation()
    slides = text.split('\n\n')
    for slide_content in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = slide_content
    prs.save(output_file)




# Example usage:
output_ppt_file = "05_Convertir_pdf_ppt.pptx"

pdf_to_ppt(output_ppt_file)